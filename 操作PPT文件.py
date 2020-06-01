from pptx import *
from pptx.chart.data import *
from pptx.util import Inches,Pt
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.dml.color import RGBColor


# 打开PPT模板
ppt = pptx.Presentation(ppt_path)
# 新建一页ppt
slide_layout = ppt.slides[0].slide_layout
ppt.slides.add_slide(slide_layout)
slide = ppt.slides[i_slide]

ppt.slides.add_slide(self.slide_layout)
self.i_slide += 1
slide = ppt.slides[self.i_slide]
# 修改车号
car_number_textbox = slide.shapes[4].text_frame.paragraphs[0]
car_number_textbox.text = Car_Number
car_number_textbox.font.size = Pt(16)
car_number_textbox.font.color.rgb = RGBColor(80, 80, 80)
# 修改图表数据
chart = slide.shapes[0].chart   # 0代表图表在PPT页面中的图形次序
chart_data = ChartData()
chart_data.categories = DayList[Start_Day_Of_Slide:Start_Day_Of_Slide + Days_In_Slide]  # x轴，即天数
chart_data.add_series(CaseLegend[iCase][0], AnalysisResult_Daily[iCase][Start_Day_Of_Slide:Start_Day_Of_Slide + Days_In_Slide])   # 抱怨数据
chart_data.add_series(CaseLegend[iCase][1], AnalysisResult_Condition_Daily[iCase][Start_Day_Of_Slide:Start_Day_Of_Slide + Days_In_Slide]) # 工况数据
chart.replace_data(chart_data)
# 获取百分比文本框的位置信息
(Percentage_Textbox_Left_First, Percentage_Textbox_Left_Last, Percentage_Textbox_Top, Percentage_Textbox_Width, Percentage_Textbox_Height)\
    = Load_Percentage_Textbox_Position_Info(Days=Days_In_Slide, yaxis_limit=max(AnalysisResult_Condition_Daily[iCase][Start_Day_Of_Slide:Start_Day_Of_Slide + Days_In_Slide]))
# 添加百分比文本框
for iDay in range(0,Days_In_Slide):
    Percentage_Textbox_Left =(Percentage_Textbox_Left_Last-Percentage_Textbox_Left_First)/(Days_In_Slide-1)*iDay+Percentage_Textbox_Left_First
    percentage_textbox = slide.shapes.add_textbox(Inches(Percentage_Textbox_Left), Inches(Percentage_Textbox_Top), Inches(Percentage_Textbox_Width), Inches(Percentage_Textbox_Height)) # 添加文本框
    percentage_textbox_paragraph = percentage_textbox.text_frame.paragraphs[0]
    percentage_textbox_paragraph_run = percentage_textbox_paragraph.add_run()
    percentage_textbox_paragraph_run.text = str(int(AnalysisResult_Percentage_Daily[iCase][iDay+Start_Day_Of_Slide] * 100)) + '%'   # 百分比计算结果
    percentage_textbox_font = percentage_textbox_paragraph_run.font
    percentage_textbox_font.bold = True
    percentage_textbox_font.size = Pt(6)


# 修改图表数据
chart = slide.shapes[0].chart
chart_data = XyChartData()
Series_Acc = chart_data.add_series(CaseLegend[iCase][0])
for iDay in range (Start_Day_Of_Slide,Start_Day_Of_Slide + Days_In_Slide):
    Series_Acc.add_data_point(Mileage_Info[iDay],AnalysisResult_Acc[iCase][iDay])
Series_ConditionAcc = chart_data.add_series(CaseLegend[iCase][1])
for iDay in range (Start_Day_Of_Slide,Start_Day_Of_Slide + Days_In_Slide):
    Series_ConditionAcc.add_data_point(Mileage_Info[iDay],AnalysisResult_Condition_Acc[iCase][iDay])
chart.replace_data(chart_data)
chart.category_axis.minimum_scale = Mileage_Info[Start_Day_Of_Slide]
chart.category_axis.maximum_scale = Mileage_Info[Start_Day_Of_Slide + Days_In_Slide-1]
# 根据抱怨和工况的差值判断并设置图表内数据的文字标签的布置方式，即是在数据点上方还是下方。
iTemp = 0
for iDay in range(Start_Day_Of_Slide, Start_Day_Of_Slide + Days_In_Slide):
    if AnalysisResult_Condition_Acc[iCase][iDay]> 0 and abs(AnalysisResult_Acc[iCase][iDay] - AnalysisResult_Condition_Acc[iCase][iDay]) / AnalysisResult_Condition_Acc[iCase][iDay] < 0.2 :
        if AnalysisResult_Acc[iCase][iDay] < AnalysisResult_Condition_Acc[iCase][iDay]:
            chart.series[0].points[iTemp].data_label.position = XL_LABEL_POSITION.BELOW
            chart.series[1].points[iTemp].data_label.position = XL_LABEL_POSITION.ABOVE
        else:
            chart.series[0].points[iTemp].data_label.position = XL_LABEL_POSITION.ABOVE
            chart.series[1].points[iTemp].data_label.position = XL_LABEL_POSITION.BELOW
    else:
        chart.series[0].points[iTemp].data_label.position = XL_LABEL_POSITION.ABOVE
        chart.series[1].points[iTemp].data_label.position = XL_LABEL_POSITION.ABOVE

slide = ppt.slides[(iCase - 1) * Slide_Number_Each_Group + 1]
chart = slide.shapes[0].chart
# 如果车辆数较多，则缩小x轴字体
if len(Car_List) > 9:
    chart.category_axis.tick_labels.font.size = Pt(10)
else:
    chart.category_axis.tick_labels.font.size = Pt(12)
chart_data = ChartData()
chart_data.categories = Car_List[1:iCar_Max + 1]  # 车号
chart_data.add_series(Case_Legend[iCase][0], Load_Analysis_result(iData=0))  # 抱怨
chart_data.add_series(Case_Legend[iCase][1], Load_Analysis_result(iData=1))  # 工况
chart.replace_data(chart_data)
# 百分比文本框位置
(Percentage_Textbox_Left_First, Percentage_Textbox_Left_Last, Percentage_Textbox_Top,
 Percentage_Textbox_Width, Percentage_Textbox_Height) = Load_Percentage_Textbox_Position_Info(
    Days=iCar_Max, yaxis_limit=max(Load_Analysis_result(iData=1)))
# 百分比文本框
for iTemp in range(0, iCar_Max):
    Percentage_Textbox_Left = (Percentage_Textbox_Left_Last - Percentage_Textbox_Left_First) / (
    iCar_Max - 1) * iTemp + Percentage_Textbox_Left_First
    percentage_textbox = slide.shapes.add_textbox(Inches(Percentage_Textbox_Left), Inches(Percentage_Textbox_Top)
                                                  , Inches(Percentage_Textbox_Width),
                                                  Inches(Percentage_Textbox_Height))
    percentage_textbox_paragraph = percentage_textbox.text_frame.paragraphs[0]
    percentage_textbox_paragraph_run = percentage_textbox_paragraph.add_run()
    percentage_textbox_paragraph_run.text = str(int(AnalysisResult[iTemp + 1][iCase][2] * 100)) + '%'
    percentage_textbox_font = percentage_textbox_paragraph_run.font
    percentage_textbox_font.bold = True
    percentage_textbox_font.size = Pt(6)

# ------------------------------------------------------------------------------------------------------------------
# 插入图片
# ------------------------------------------------------------------------------------------------------------------
slide = ppt.slides[1 + n_car]
textbox_date = slide.shapes[0].text_frame.paragraphs[0]
textbox_date.text = folder_name
slide.shapes.add_picture(image_file=pic_path_speed, left=Inches(0.41), top=Inches(1.2), width=Inches(9.45),
                         height=Inches(1.6))
slide.shapes.add_picture(image_file=pic_path_engine_speed, left=Inches(0.33), top=Inches(3), width=Inches(8.83),
                         height=Inches(1.6))
slide.shapes.add_picture(image_file=pic_path_pedal, left=Inches(0.41), top=Inches(4.8), width=Inches(8.75),
                         height=Inches(1.6))

# 填写表格数据
slide = ppt.slides[0]
textbox_date = slide.shapes[0].text_frame.paragraphs[0]
textbox_date.text = 'SVW OBD Driving Behavior Daily Summary'
table = slide.shapes[1].table
for i in range(len(vehicle_number)):
    # 写入表格内容
    table.cell(2 + i, 0).text = str(vehicle_number[i])
    table.cell(2 + i, 1).text_frame.text = str(project_mileage[i])
    table.cell(2 + i, 2).text_frame.text = str(data_record_date[i])
    table.cell(2 + i, 3).text_frame.text = str(drive_time[i])
    table.cell(2 + i, 4).text_frame.text = str(speed_max[i])
    table.cell(2 + i, 5).text_frame.text = '{0} & {1}'.format(speed_above_80[i], speed_above_120[i])
    table.cell(2 + i, 6).text_frame.text = str(speed_average[i])
    table.cell(2 + i, 7).text_frame.text = str(engine_speed_max[i])
    table.cell(2 + i, 8).text_frame.text = str(engine_speed_above_3800[i])
    table.cell(2 + i, 9).text_frame.text = str(pedal_max[i])
    table.cell(2 + i, 10).text_frame.text = str(pedal_above_70[i])
    # 修正文字大小
    table.cell(2 + i, 0).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 1).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 2).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 3).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 4).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 5).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 6).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 7).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 8).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 9).text_frame.paragraphs[0].font.size = Pt(10)
    table.cell(2 + i, 10).text_frame.paragraphs[0].font.size = Pt(10)

# ------------------------------------------------------------------------------------------------------------------

# 保存PPT
ppt.save(Result_Path + "\\" + Car_Number + ".pptx")